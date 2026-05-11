#!/usr/bin/env python3
"""convert_pptx.py — Convert any PPTX to ADATA template style.

Only replaces the **background**, **title**, and **subtitle** styling.
All other content (text boxes, body text, images, tables, shapes) is
preserved exactly as-is from the source deck.

Usage:
    python .agents/skills/adata-pptx/scripts/convert_pptx.py  source.pptx  [output.pptx]

Slide auto-classification (by position and content):
  slide 0              → Cover   (ADATA cover layout background)
  title-only slides    → Section Divider  (ADATA divider background)
  content slides       → Content slide   (ADATA content background)
  last slide (empty)   → Closing/blank

Section colour cycle:  Blue → Green → Orange → Magenta

Dependencies:  pip install python-pptx lxml
"""

import argparse
import copy
import io
import random
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches
from lxml import etree

# ─── Brand constants ──────────────────────────────────────────────────────────

SECTION_COLORS = [
    RGBColor(0x50, 0x97, 0xFF),  # Section 1 – Blue    #5097FF
    RGBColor(0x19, 0xC7, 0x11),  # Section 2 – Green   #19C711
    RGBColor(0xFF, 0x90, 0x00),  # Section 3 – Orange  #FF9000
    RGBColor(0xFF, 0x47, 0xFF),  # Section 4 – Magenta #FF47FF
]
DARK_NAVY = RGBColor(0x0E, 0x28, 0x41)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SKILL_DIR     = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = SKILL_DIR / "adata-template.pptx"

_POTX_CT_OLD = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
_POTX_CT_NEW = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"


def _load_template(template_path: Path) -> Presentation:
    """
    Load the ADATA template even if it still carries the .potx content-type.
    If python-pptx rejects it, we patch [Content_Types].xml in-memory and
    retry from a temporary file.
    """
    try:
        return Presentation(template_path)
    except ValueError:
        # Template is a .potx — patch the content type in memory
        buf = io.BytesIO()
        with zipfile.ZipFile(template_path, "r") as zin, \
             zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        _POTX_CT_OLD.encode(),
                        _POTX_CT_NEW.encode(),
                    )
                zout.writestr(item, data)
        buf.seek(0)
        return Presentation(buf)

# Layout indices inside adata-template.pptx (0-based)
# Template slide mapping:
#   Slide  1 → layout 0  (標題投影片)    Cover
#   Slide  2 → layout 1  (標題及內容)    Agenda
#   Slide  3 → layout 2  (章節標題)      Section 1 divider
#   Slide  4 → layout 2  (章節標題)      Section 1 content  ← same layout as divider
#   Slide  5 → layout 4  (比較)          Section 2 divider
#   Slide  6 → layout 4  (比較)          Section 2 content  ← same layout as divider
#   Slide  7 → layout 6  (1_比較)        Section 3 divider
#   Slide  8 → layout 7  (1_只有標題)    Section 3 content
#   Slide  9 → layout 8  (2_比較)        Section 4 divider
#   Slide 10 → layout 9  (2_只有標題)    Section 4 content
#   Slide 11 → layout 10 (空白)          Closing
_LO_COVER  = 0   # 標題投影片  — Cover
_LO_BLANK  = 10  # 空白        — Closing slide

# Paired (divider_layout, content_layout) for each of the 4 ADATA section styles.
# Randomly pick one pair per section; divider uses divider_layout, all content
# slides in that section use content_layout.
# Content layouts must have title at top ≈ 0.09" (content position, not section position).
# Layouts 3/5/7/9 all have title at top=0.09"; layouts 2/4/6/8 have title at top=1.31" (section).
_SECTION_PAIRS = [
    (2, 3),   # 章節標題(divider) / 兩個內容(content)  — Section 1 colour style
    (4, 5),   # 比較(divider)     / 只有標題(content)  — Section 2 colour style
    (6, 7),   # 1_比較(divider)   / 1_只有標題(content) — Section 3 colour style
    (8, 9),   # 2_比較(divider)   / 2_只有標題(content) — Section 4 colour style
]
_PAIR_NAMES = [
    ('章節標題', '兩個內容'),
    ('比較', '只有標題'),
    ('1_比較', '1_只有標題'),
    ('2_比較', '2_只有標題'),
]


# ─── Content extraction ───────────────────────────────────────────────────────

class SlideData:
    """All extracted content from one source slide."""
    __slots__ = ("title", "subtitle", "body", "images", "tables", "notes")

    def __init__(self):
        self.title:    str = ""
        self.subtitle: str = ""
        # Each body item: {"text": str, "level": int,
        #   "bullet": "char"|"auto"|"none"|None, "char": str|None,
        #   "auto_type": str|None}
        self.body:   list[dict] = []
        # Each image: {blob, ext, left, top, width, height}
        self.images: list[dict] = []
        # Each table: a 2-D list of strings  [row][col]
        self.tables: list[list[list[str]]] = []
        self.notes:  str = ""


def _max_font_size(shape) -> int:
    """Return the largest font size (in hundredths of a point) found in *shape*."""
    best = 0
    if not shape.has_text_frame:
        return 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is not None:
                sz = rPr.get('sz')
                if sz:
                    best = max(best, int(sz))
    return best


def _extract_para(para) -> dict:
    """Extract a single paragraph's text + bullet/indent metadata."""
    item = {"text": para.text.strip(), "level": 0,
            "bullet": None, "char": None, "auto_type": None}
    pPr = para._p.find(qn('a:pPr'))
    if pPr is not None:
        item["level"] = int(pPr.get('lvl', '0'))
        buChar = pPr.find(qn('a:buChar'))
        buAuto = pPr.find(qn('a:buAutoNum'))
        buNone = pPr.find(qn('a:buNone'))
        if buChar is not None:
            item["bullet"] = "char"
            item["char"] = buChar.get('char', '•')
        elif buAuto is not None:
            item["bullet"] = "auto"
            item["auto_type"] = buAuto.get('type', 'arabicPeriod')
        elif buNone is not None:
            item["bullet"] = "none"
    return item


def _clean_xml_string(text: str) -> str:
    """Remove XML-incompatible characters (NULL bytes, control characters)."""
    if not text:
        return text
    # Remove NULL bytes and other XML-incompatible control characters
    # Keep common valid characters: alphanumerics, spaces, punctuation, CJK, emoji
    cleaned = ''.join(c for c in text if ord(c) >= 32 or c in '\t\n\r')
    return cleaned


def _plain_body(text: str) -> dict:
    """Create a plain body item (no bullet, level 0)."""
    return {"text": text, "level": 0, "bullet": None, "char": None, "auto_type": None}


def _extract_shapes(shapes, data: SlideData, is_title_done: bool) -> bool:
    """
    Walk shape tree, fill *data* in-place.
    Returns True once the title placeholder has been consumed.
    Call recursively for group shapes.
    """
    for shape in shapes:
        # ── Grouped shapes (recurse) ─────────────────────────────────────
        try:
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                is_title_done = _extract_shapes(shape.shapes, data, is_title_done)
                continue
        except Exception:
            pass

        # ── Placeholders ─────────────────────────────────────────────────
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            if idx == 0 and not is_title_done and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    data.title = t
                    is_title_done = True
            elif idx == 1 and shape.has_text_frame:
                # Subtitle placeholder
                t = shape.text_frame.text.strip()
                if t and not data.subtitle:
                    data.subtitle = t
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    item = _extract_para(para)
                    if item["text"]:
                        data.body.append(item)
            continue

        # ── Pictures ─────────────────────────────────────────────────────
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                data.images.append({
                    "blob":   shape.image.blob,
                    "ext":    shape.image.ext,
                    "left":   shape.left,
                    "top":    shape.top,
                    "width":  shape.width,
                    "height": shape.height,
                    "name":   shape.name,
                })
            except Exception as exc:
                print(f"  WARNING: skipped image '{getattr(shape, 'name', '?')}': {exc}",
                      file=sys.stderr)
            continue

        # ── Tables ───────────────────────────────────────────────────────
        if shape.has_table:
            rows = [
                [cell.text_frame.text for cell in row.cells]
                for row in shape.table.rows
            ]
            data.tables.append(rows)
            continue

        # ── Free text boxes / auto-shapes ─────────────────────────────────
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                item = _extract_para(para)
                if item["text"]:
                    data.body.append(item)

    return is_title_done


def extract_slide(slide) -> SlideData:
    data = SlideData()
    _extract_shapes(slide.shapes, data, is_title_done=False)

    # ── Title recovery: if no placeholder title was found, scan all shapes
    #    and pick the one with the LARGEST font as the title. ──
    has_ph_title = any(
        s.is_placeholder and s.placeholder_format.idx == 0
        and s.has_text_frame and s.text_frame.text.strip()
        for s in slide.shapes
    )
    if not has_ph_title and data.body:
        best_text = ""
        best_sz   = 0
        for shape in slide.shapes:
            if shape.is_placeholder or shape.shape_type == 13 or shape.has_table:
                continue
            if not shape.has_text_frame:
                continue
            sz = _max_font_size(shape)
            if sz > best_sz:
                best_sz   = sz
                best_text = shape.text_frame.text.strip()
        if best_text and best_sz > 0:
            data.title = best_text
            # Remove title lines from body
            title_lines = {l.strip() for l in best_text.replace('\r\n', '\n').split('\n') if l.strip()}
            data.body = [item for item in data.body if item["text"] not in title_lines]
        elif not data.title and data.body:
            # Fallback: first body item becomes title
            data.title = data.body.pop(0)["text"]

    # Speaker notes
    if slide.has_notes_slide:
        try:
            data.notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

    return data


def classify(data: SlideData, idx: int, total: int) -> str:
    """Return ADATA slide type: 'cover' | 'divider' | 'content' | 'closing'."""
    if idx == 0:
        return "cover"
    if idx == total - 1 and not data.body and not data.images and not data.tables:
        return "closing"
    # Treat title-only (or very sparse) slides as section dividers
    if not data.body and not data.images and not data.tables:
        return "divider"
    if (data.title
            and len(data.body) <= 2
            and not data.images
            and not data.tables):
        return "divider"
    return "content"


# ─── Template management ──────────────────────────────────────────────────────

def delete_slide(prs: Presentation, index: int) -> None:
    """Remove the slide at *index* from *prs* (keeps master/layouts)."""
    sldIdLst = prs.slides._sldIdLst
    sldId    = sldIdLst[index]
    rId      = sldId.get(qn("r:id"))
    sldIdLst.remove(sldId)
    try:
        prs.part.drop_rel(rId)
    except KeyError:
        pass


def clear_all_slides(prs: Presentation) -> None:
    """Remove every slide, leaving masters and layouts intact."""
    while len(prs.slides):
        delete_slide(prs, 0)


# ─── Helpers: XML-based text styling ─────────────────────────────────────────

def _is_cjk(text: str) -> bool:
    return any('\u2E80' <= c <= '\u9FFF' or '\uF900' <= c <= '\uFAFF' for c in text)


def _clear_paragraphs(txBody) -> None:
    """Remove all <a:p> elements, preserve lstStyle and other metadata."""
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)


def _make_styled_run(text: str, color: RGBColor, size_pt: int | None = None) -> etree._Element:
    """Build <a:r> with Arial Black, explicit color, correct lang.
    If size_pt is None, no sz attribute is written so the template placeholder
    size is inherited automatically.
    """
    text = _clean_xml_string(text)
    lang = 'zh-TW' if _is_cjk(text) else 'en-US'
    alt  = 'en-US' if lang == 'zh-TW' else 'zh-TW'
    color_hex = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"

    r    = etree.Element(qn('a:r'))
    rPr  = etree.SubElement(r, qn('a:rPr'))
    rPr.set('lang', lang)
    rPr.set('altLang', alt)
    if size_pt is not None:
        rPr.set('sz', str(int(size_pt * 100)))
    rPr.set('dirty', '0')
    sf   = etree.SubElement(rPr, qn('a:solidFill'))
    clr  = etree.SubElement(sf, qn('a:srgbClr'))
    clr.set('val', color_hex)
    lat  = etree.SubElement(rPr, qn('a:latin'))
    lat.set('typeface', 'Arial Black')
    t    = etree.SubElement(r, qn('a:t'))
    t.text = text
    return r


def _make_plain_run(text: str) -> etree._Element:
    """Build <a:r> with minimal rPr so it inherits layout formatting."""
    text = _clean_xml_string(text)
    lang = 'zh-TW' if _is_cjk(text) else 'en-US'
    alt  = 'en-US' if lang == 'zh-TW' else 'zh-TW'
    r    = etree.Element(qn('a:r'))
    rPr  = etree.SubElement(r, qn('a:rPr'))
    rPr.set('lang', lang)
    rPr.set('altLang', alt)
    rPr.set('dirty', '0')
    t    = etree.SubElement(r, qn('a:t'))
    t.text = text
    return r


def _set_title(tf, text: str, color: RGBColor, size_pt: int | None = None) -> None:
    """
    Set title text with Arial Black + color via direct XML.
    Splits on \\n into separate paragraphs (renders correctly in PowerPoint).
    If size_pt is None, the template placeholder's font size is inherited.
    """
    txBody = tf._txBody
    _clear_paragraphs(txBody)
    lines = [l for l in text.replace('\r\n', '\n').replace('\r', '\n').split('\n') if l.strip()]
    if not lines:
        lines = [text or ' ']
    for line in lines:
        p   = etree.SubElement(txBody, qn('a:p'))
        p.append(_make_styled_run(line, color, size_pt))
        end = etree.SubElement(p, qn('a:endParaRPr'))
        end.set('lang', 'zh-TW' if _is_cjk(line) else 'en-US')
        end.set('dirty', '0')


def _fill_ph_text(tf, text: str, color: RGBColor | None = None) -> None:
    """Set placeholder text, inheriting ALL template formatting.

    Only the text content (and optionally the font colour) is changed;
    font family, size, bold, position, etc. all come from the template
    placeholder / layout / master.
    """
    text = _clean_xml_string(text)
    txBody = tf._txBody
    _clear_paragraphs(txBody)
    lines = [l for l in text.replace('\r\n', '\n').replace('\r', '\n').split('\n') if l.strip()]
    if not lines:
        lines = [text or ' ']
    for line in lines:
        line = _clean_xml_string(line)
        lang = 'zh-TW' if _is_cjk(line) else 'en-US'
        alt  = 'en-US' if lang == 'zh-TW' else 'zh-TW'

        p   = etree.SubElement(txBody, qn('a:p'))
        r   = etree.SubElement(p, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'))
        rPr.set('lang', lang)
        rPr.set('altLang', alt)
        rPr.set('dirty', '0')
        if color is not None:
            sf  = etree.SubElement(rPr, qn('a:solidFill'))
            clr = etree.SubElement(sf, qn('a:srgbClr'))
            clr.set('val', f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
        t = etree.SubElement(r, qn('a:t'))
        t.text = line
        end = etree.SubElement(p, qn('a:endParaRPr'))
        end.set('lang', lang)
        end.set('dirty', '0')


def _set_body_text(tf, paragraphs: list) -> None:
    """
    Fill body placeholder with paragraphs that preserve bullet/indent info.
    Each item is either a dict (from _extract_para) or a plain str.
    """
    txBody = tf._txBody
    _clear_paragraphs(txBody)
    if not paragraphs:
        p = etree.SubElement(txBody, qn('a:p'))
        etree.SubElement(p, qn('a:endParaRPr')).set('dirty', '0')
        return
    for item in paragraphs:
        # Accept both str and dict
        if isinstance(item, str):
            item = {"text": item, "level": 0, "bullet": None, "char": None, "auto_type": None}
        text  = item["text"]
        level = item.get("level", 0)

        p   = etree.SubElement(txBody, qn('a:p'))
        # ── Paragraph properties: level + bullet ──
        pPr = etree.SubElement(p, qn('a:pPr'))
        if level > 0:
            pPr.set('lvl', str(level))
        bullet = item.get("bullet")
        if bullet == "char":
            bu = etree.SubElement(pPr, qn('a:buChar'))
            bu.set('char', item.get("char") or '•')
        elif bullet == "auto":
            bu = etree.SubElement(pPr, qn('a:buAutoNum'))
            bu.set('type', item.get("auto_type") or 'arabicPeriod')
        elif bullet == "none":
            etree.SubElement(pPr, qn('a:buNone'))
        # If bullet is None, no bullet element → inherits from layout lstStyle

        p.append(_make_plain_run(text))
        end = etree.SubElement(p, qn('a:endParaRPr'))
        end.set('lang', 'zh-TW' if _is_cjk(text) else 'en-US')
        end.set('dirty', '0')


# ─── Helpers: add images & tables ────────────────────────────────────────────

def _scale_coord(value: int, src_size: int, dst_size: int) -> int:
    if src_size == 0:
        return value
    return int(value * dst_size / src_size)


def add_image(slide, img: dict,
              src_w: int, src_h: int,
              dst_w: int, dst_h: int) -> None:
    """Add an image to *slide*, scaling coordinates from source to dest slide size."""
    left   = _scale_coord(img["left"],   src_w, dst_w)
    top    = _scale_coord(img["top"],    src_h, dst_h)
    width  = _scale_coord(img["width"],  src_w, dst_w)
    height = _scale_coord(img["height"], src_h, dst_h)

    # Clamp to slide bounds
    left   = max(0, min(left,   dst_w - 914400))   # at least 1 inch from right edge
    top    = max(0, min(top,    dst_h - 914400))
    width  = max(914400, min(width,  dst_w - left))  # at least 1 inch wide
    height = max(914400, min(height, dst_h - top))

    slide.shapes.add_picture(io.BytesIO(img["blob"]), left, top, width, height)


def add_table(slide, tbl_data: list[list[str]],
              left: int, top: int, width: int, height: int,
              accent: RGBColor) -> None:
    """
    Add *tbl_data* as a table shape.
    The first row receives a dark-navy header fill with white bold text.
    Subsequent rows use the section accent colour for the first cell.
    """
    if not tbl_data:
        return

    rows = len(tbl_data)
    cols = max(len(r) for r in tbl_data)
    if cols == 0:
        return

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl       = tbl_shape.table

    for ri, row_data in enumerate(tbl_data):
        for ci in range(cols):
            cell_text = row_data[ci] if ci < len(row_data) else ""
            cell = tbl.cell(ri, ci)
            cell.text = cell_text

            if ri == 0:
                # Header row: dark-navy fill, white bold text
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold  = True
                        run.font.color.rgb = WHITE
                    para.font.bold  = True
                    para.font.color.rgb = WHITE
                # Apply solid fill via OOXML
                tcPr = cell._tc.get_or_add_tcPr()
                # Remove any existing fill
                for child in list(tcPr):
                    if child.tag.endswith(("solidFill", "gradFill", "noFill", "pattFill")):
                        tcPr.remove(child)
                sf = etree.SubElement(tcPr, qn("a:solidFill"))
                etree.SubElement(sf, qn("a:srgbClr")).set("val", "0E2841")

            elif ci == 0:
                # First column of data rows: accent colour text
                for para in cell.text_frame.paragraphs:
                    para.font.bold = True
                    para.font.color.rgb = accent


# ─── Slide builders ───────────────────────────────────────────────────────────

def _placeholders(slide) -> dict:
    """Return {idx: placeholder} dict for *slide*."""
    return {ph.placeholder_format.idx: ph for ph in slide.placeholders}


def _find_body_ph(phs: dict):
    """Find the body-content placeholder (varies by layout: idx 2, 13, 10, …)."""
    for idx in (2, 13, 10):
        if idx in phs:
            return phs[idx]
    return None


def _find_subtitle_ph(phs: dict):
    """Find the subtitle placeholder (idx 1, 10, 11, …)."""
    for idx in (1, 10, 11):
        if idx in phs:
            return phs[idx]
    return None


def build_cover(prs_out: Presentation, data: SlideData) -> None:
    layout = prs_out.slide_layouts[_LO_COVER]
    slide  = prs_out.slides.add_slide(layout)
    phs    = _placeholders(slide)

    # Cover title: white, Arial Black — size inherited from template (66pt)
    if 0 in phs and data.title:
        _set_title(phs[0].text_frame, data.title, WHITE)

    # Subtitle — cover PH indices are typically 10, 11 (not 1)
    sub_ph = _find_subtitle_ph(phs)
    if sub_ph:
        subtitle_parts = []
        if data.subtitle:
            subtitle_parts.append(data.subtitle)
        subtitle_parts.extend(item["text"] for item in data.body[:3])
        if subtitle_parts:
            _set_body_text(sub_ph.text_frame, subtitle_parts[:3])

    # Images on cover are placed as-is
    sw, sh = int(prs_out.slide_width), int(prs_out.slide_height)
    for img in data.images:
        try:
            add_image(slide, img, sw, sh, sw, sh)
        except Exception as exc:
            print(f"  WARNING: cover image skipped: {exc}", file=sys.stderr)


def build_divider(prs_out: Presentation, data: SlideData,
                  section_idx: int, lo_idx: int = 2) -> None:
    layout = prs_out.slide_layouts[lo_idx]
    slide  = prs_out.slides.add_slide(layout)
    phs    = _placeholders(slide)

    if 0 in phs and data.title:
        _set_title(phs[0].text_frame, data.title, WHITE)

    # Subtitles — divider uses idx=10
    body_ph = _find_body_ph(phs)
    if body_ph and data.body:
        # Divider subtitles: flatten to plain text (no bullets)
        plain = [item["text"] for item in data.body[:3]]
        _set_body_text(body_ph.text_frame, plain)

    _copy_notes(slide, data.notes)


def build_content(prs_out: Presentation, data: SlideData,
                  section_idx: int,
                  src_w: int, src_h: int) -> None:
    color  = SECTION_COLORS[min(section_idx, 3)]
    dst_w  = int(prs_out.slide_width)
    dst_h  = int(prs_out.slide_height)

    # Always use "只有標題" (Title Only) — it has PH idx=0 (title),
    # idx=1 (subtitle), idx=2 (body) with proper template styling.
    layout = prs_out.slide_layouts[_LO_TITLE_ONLY]
    slide  = prs_out.slides.add_slide(layout)
    phs    = _placeholders(slide)

    # ── Title (styled: Arial Black + section accent colour) — size inherited from template (55pt)
    if 0 in phs and data.title:
        _set_title(phs[0].text_frame, data.title, color)

    # ── Subtitle (inherits layout styling) ──
    sub_ph = _find_subtitle_ph(phs)
    if sub_ph and data.subtitle:
        _set_body_text(sub_ph.text_frame, [data.subtitle])

    # ── Body (inherits layout styling — font, size, bullets) ──
    body_ph = _find_body_ph(phs)
    if body_ph and data.body:
        _set_body_text(body_ph.text_frame, data.body)

    # ── Images (preserve original positions, scaled to ADATA slide) ──
    for img in data.images:
        try:
            add_image(slide, img, src_w, src_h, dst_w, dst_h)
        except Exception as exc:
            print(f"  WARNING: image skipped: {exc}", file=sys.stderr)

    # ── Tables ──
    if data.tables:
        margin    = int(Inches(0.5))
        tbl_left  = margin
        tbl_width = dst_w - 2 * margin
        tbl_top   = int(Inches(1.8))

        for tbl_data in data.tables:
            rows       = len(tbl_data)
            row_height = int(Inches(0.4))
            tbl_height = max(rows * row_height, int(Inches(1.0)))
            if tbl_top + tbl_height > dst_h - margin:
                tbl_height = dst_h - tbl_top - margin
            if tbl_height <= 0:
                print("  WARNING: table clipped off slide, skipping.", file=sys.stderr)
                break
            try:
                add_table(slide, tbl_data, tbl_left, tbl_top,
                          tbl_width, tbl_height, color)
            except Exception as exc:
                print(f"  WARNING: table skipped: {exc}", file=sys.stderr)
            tbl_top += tbl_height + int(Inches(0.2))

    _copy_notes(slide, data.notes)


def build_closing(prs_out: Presentation, data: SlideData) -> None:
    layout = prs_out.slide_layouts[_LO_BLANK]
    slide  = prs_out.slides.add_slide(layout)
    # Add title text as a centred text box if present
    if data.title or data.body:
        text = data.title or (data.body[0]["text"] if isinstance(data.body[0], dict) else data.body[0])
        sw   = int(prs_out.slide_width)
        sh   = int(prs_out.slide_height)
        txBox = slide.shapes.add_textbox(
            int(Inches(1)), int(sh * 0.35),
            sw - int(Inches(2)), int(Inches(1.5))
        )
        _set_title(txBox.text_frame, text, WHITE, 32)
    _copy_notes(slide, data.notes)


def _copy_notes(slide, notes: str) -> None:
    if not notes:
        return
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except Exception:
        pass


# ─── Shape cloning helpers ────────────────────────────────────────────────────

# Relationship types that belong to the slide layout/master — skip when cloning
_SKIP_REL_TYPES = {
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout',
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster',
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster',
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide',
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags',
}

_R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def _get_ph_idx(sp_element) -> int | None:
    """Return the placeholder idx from a p:sp element, or None."""
    nvSpPr = sp_element.find(qn('p:nvSpPr'))
    if nvSpPr is None:
        return None
    nvPr = nvSpPr.find(qn('p:nvPr'))
    if nvPr is None:
        return None
    ph = nvPr.find(qn('p:ph'))
    if ph is None:
        return None
    idx = ph.get('idx')
    return int(idx) if idx is not None else 0  # idx defaults to 0 (title)


def _remove_ph_shapes(slide, *, keep_idx: set[int] | None = None) -> None:
    """Remove placeholder shapes from slide's spTree.

    If *keep_idx* is given, placeholders with those idx values are kept;
    all others are removed.  If *keep_idx* is None, ALL placeholders are removed.
    """
    spTree = slide.shapes._spTree
    to_remove = []
    for child in list(spTree):
        if child.tag == qn('p:sp'):
            ph_idx = _get_ph_idx(child)
            if ph_idx is not None:
                if keep_idx is not None and ph_idx in keep_idx:
                    continue  # keep this one
                to_remove.append(child)
    for el in to_remove:
        spTree.remove(el)


def _build_rId_map(src_slide, dst_slide) -> dict:
    """Copy ALL content relationships from source slide to dest slide.

    Returns a mapping {old_rId: new_rId} so cloned XML can be patched.
    Skips layout/master rels (the dest already has ADATA ones).
    """
    rId_map = {}
    for rId in list(src_slide.part.rels.keys()):
        rel = src_slide.part.rels[rId]
        if rel.reltype in _SKIP_REL_TYPES:
            continue
        try:
            if rel.is_external:
                new_rId = dst_slide.part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            else:
                new_rId = dst_slide.part.relate_to(
                    rel.target_part, rel.reltype
                )
            rId_map[rId] = new_rId
        except Exception as e:
            print(f"  WARNING: could not copy relationship {rId} "
                  f"({rel.reltype.split('/')[-1]}): {e}", file=sys.stderr)
    return rId_map


def _replace_rIds(element, rId_map: dict) -> None:
    """Replace ALL relationship ID references in an XML element tree.

    Walks every element and checks every attribute in the ``r:`` namespace
    (r:embed, r:link, r:id, etc.) as well as common non-namespaced rId attrs.
    """
    r_prefix = '{' + _R_NS + '}'
    for el in element.iter():
        for attr_name in list(el.attrib.keys()):
            if attr_name.startswith(r_prefix):
                old_val = el.get(attr_name)
                if old_val in rId_map:
                    el.set(attr_name, rId_map[old_val])


def _scale_xfrm(xfrm, src_w, src_h, dst_w, dst_h) -> None:
    """Scale an a:xfrm element's offset and extent."""
    off = xfrm.find(qn('a:off'))
    if off is not None:
        x = int(off.get('x', '0'))
        y = int(off.get('y', '0'))
        off.set('x', str(int(x * dst_w / src_w)))
        off.set('y', str(int(y * dst_h / src_h)))
    ext = xfrm.find(qn('a:ext'))
    if ext is not None:
        cx = int(ext.get('cx', '0'))
        cy = int(ext.get('cy', '0'))
        ext.set('cx', str(int(cx * dst_w / src_w)))
        ext.set('cy', str(int(cy * dst_h / src_h)))


def _clone_slide_shapes(src_slide, dst_slide,
                        src_w: int, src_h: int,
                        dst_w: int, dst_h: int,
                        skip_ph_idx: set[int] | None = None) -> None:
    """Clone shapes from *src_slide* onto *dst_slide*.

    If *skip_ph_idx* is given, source placeholders with those idx values are
    NOT cloned (the dest template's own placeholders are used instead).

    1. Copy every relationship (images, charts, OLE, hyperlinks…) from source
       to dest, building a complete rId mapping.
    2. Deep-copy each shape element's XML.
    3. Patch all ``r:*`` attributes using the rId mapping.
    4. Scale coordinates if slide dimensions differ.
    """
    # Step 1 — Copy relationships
    rId_map = _build_rId_map(src_slide, dst_slide)

    # Step 2–4 — Clone shapes
    src_spTree = src_slide.shapes._spTree
    dst_spTree = dst_slide.shapes._spTree
    need_scale = (src_w != dst_w or src_h != dst_h)

    # Only copy actual shape elements (sp, pic, graphicFrame, grpSp, cxnSp)
    shape_tags = {qn('p:sp'), qn('p:pic'), qn('p:graphicFrame'),
                  qn('p:grpSp'), qn('p:cxnSp')}

    for child in list(src_spTree):
        if child.tag not in shape_tags:
            continue

        # Skip source placeholders that the template already provides
        if skip_ph_idx and child.tag == qn('p:sp'):
            ph_idx = _get_ph_idx(child)
            if ph_idx is not None and ph_idx in skip_ph_idx:
                continue

        el = copy.deepcopy(child)

        # Patch relationship references
        _replace_rIds(el, rId_map)

        # Scale coordinates if slide sizes differ
        if need_scale:
            # p:sp and p:cxnSp → p:spPr/a:xfrm
            spPr = el.find(qn('p:spPr'))
            if spPr is not None:
                xfrm = spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    _scale_xfrm(xfrm, src_w, src_h, dst_w, dst_h)
            # p:grpSp → p:grpSpPr/a:xfrm  (group outer bounds)
            grpSpPr = el.find(qn('p:grpSpPr'))
            if grpSpPr is not None:
                xfrm = grpSpPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    _scale_xfrm(xfrm, src_w, src_h, dst_w, dst_h)
            # p:pic → p:spPr/a:xfrm
            pic_spPr = el.find(qn('p:spPr'))
            if pic_spPr is not None and el.tag == qn('p:pic'):
                xfrm = pic_spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    _scale_xfrm(xfrm, src_w, src_h, dst_w, dst_h)
            # p:graphicFrame → p:xfrm
            gf_xfrm = el.find(qn('p:xfrm'))
            if gf_xfrm is not None:
                _scale_xfrm(gf_xfrm, src_w, src_h, dst_w, dst_h)

        dst_spTree.append(el)


def _restyle_title_ph(slide, color: RGBColor) -> None:
    """Find the title placeholder (idx 0) and restyle with ADATA branding."""
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.idx != 0:
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text:
            _set_title(shape.text_frame, text, color)
        return


def _restyle_subtitle_ph(slide, color: RGBColor) -> None:
    """Find the subtitle placeholder and restyle with ADATA branding."""
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.idx not in (1, 10, 11):
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text:
            _set_title(shape.text_frame, text, color, size_pt=28)
        return


# ─── Main conversion ──────────────────────────────────────────────────────────

def convert(src: Path, dst: Path) -> None:
    """Convert *src* to ADATA style.

    Strategy: create ADATA-template slides (for background), then clone ALL
    source shapes onto them. Only title and subtitle placeholders are restyled;
    every other element (text boxes, body text, images, tables, charts) is
    preserved as-is from the source.
    """
    print(f"Source : {src}")
    print(f"Output : {dst}")

    # ── Load source ─────────────────────────────────────────────────────────
    src_prs = Presentation(src)
    src_w   = int(src_prs.slide_width)
    src_h   = int(src_prs.slide_height)
    total   = len(src_prs.slides)

    # ── Classify slides ─────────────────────────────────────────────────────
    print(f"\nClassifying {total} slides …")
    slides_info: list[tuple[str, SlideData, int]] = []
    section_idx = -1
    for i, slide in enumerate(src_prs.slides):
        data = extract_slide(slide)
        kind = classify(data, i, total)
        if kind == 'divider':
            section_idx += 1
        sec = max(section_idx, 0)
        slides_info.append((kind, data, sec))

        extras = []
        if data.images: extras.append(f"{len(data.images)} img")
        if data.tables: extras.append(f"{len(data.tables)} tbl")
        label = data.title[:60] or "(no title)"
        print(f"  [{i+1:>2}/{total}] {kind:<10}  {label}"
              + (f"  [{', '.join(extras)}]" if extras else ""))

    # ── Build ADATA output ──────────────────────────────────────────────────
    print(f"\nBuilding ADATA deck (background + title/subtitle only) …")
    out_prs = _load_template(TEMPLATE_PATH)
    clear_all_slides(out_prs)
    dst_w = int(out_prs.slide_width)
    dst_h = int(out_prs.slide_height)

    # Pre-assign a random (divider, content) layout pair to each section
    total_sections = max(s[2] for s in slides_info) + 1
    section_layout_map = {
        sec: random.choice(list(zip(_SECTION_PAIRS, _PAIR_NAMES)))
        for sec in range(total_sections)
    }
    # section_layout_map[sec] = ((div_lo, con_lo), (div_name, con_name))
    print("  Section layout assignments:")
    for sec, ((div_lo, con_lo), (div_name, con_name)) in section_layout_map.items():
        print(f"    Section {sec+1}: divider={div_name}(layout {div_lo})  content={con_name}(layout {con_lo})")

    for i, src_slide in enumerate(src_prs.slides):
        kind, data, sec_idx = slides_info[i]
        label = data.title[:50] or "(no title)"

        # ── Cover / Divider / Closing: use full ADATA template styling ──
        if kind == 'cover':
            build_cover(out_prs, data)
            color_name = ["Blue", "Green", "Orange", "Magenta"][min(sec_idx, 3)]
            print(f"  [{i+1:>2}/{total}] {kind:<10}  {label}  "
                  f"(section {sec_idx+1} / {color_name})")
            continue

        if kind == 'divider':
            (div_lo, _con_lo), _ = section_layout_map[sec_idx]
            build_divider(out_prs, data, sec_idx, lo_idx=div_lo)
            color_name = ["Blue", "Green", "Orange", "Magenta"][min(sec_idx, 3)]
            print(f"  [{i+1:>2}/{total}] {kind:<10}  {label}  "
                  f"(section {sec_idx+1} / {color_name})")
            continue

        if kind == 'closing':
            build_closing(out_prs, data)
            color_name = ["Blue", "Green", "Orange", "Magenta"][min(sec_idx, 3)]
            print(f"  [{i+1:>2}/{total}] {kind:<10}  {label}  "
                  f"(section {sec_idx+1} / {color_name})")
            continue

        # ── Content slides: ADATA background + template title/subtitle ──
        (_div_lo, lo_idx), _ = section_layout_map[sec_idx]
        layout    = out_prs.slide_layouts[lo_idx]
        dst_slide = out_prs.slides.add_slide(layout)

        # Keep template's title (idx 0) and subtitle (idx 1) placeholders;
        # remove any other template placeholders (body etc.)
        _remove_ph_shapes(dst_slide, keep_idx={0, 1})

        # Clone source shapes EXCEPT title/subtitle placeholders
        # (those are handled by the template placeholders above)
        _clone_slide_shapes(src_slide, dst_slide, src_w, src_h, dst_w, dst_h,
                            skip_ph_idx={0, 1})

        # ── Fill template title placeholder with source text ──
        # No color override — fully inherits template placeholder styling
        phs = _placeholders(dst_slide)
        if 0 in phs and data.title:
            _fill_ph_text(phs[0].text_frame, data.title)

        # ── Fill template subtitle placeholder with source text ──
        sub_ph = _find_subtitle_ph(phs)
        if sub_ph and data.subtitle:
            _fill_ph_text(sub_ph.text_frame, data.subtitle)

        # ── Speaker notes ──
        _copy_notes(dst_slide, data.notes)

        color_name = ["Blue", "Green", "Orange", "Magenta"][min(sec_idx, 3)]
        print(f"  [{i+1:>2}/{total}] {kind:<10}  {label}  "
              f"(section {sec_idx+1} / {color_name})")

    # ── Save ────────────────────────────────────────────────────────────────
    print(f"\nSaving → {dst}")
    out_prs.save(dst)
    print("Done.  Open in PowerPoint and review.")


# ─── CLI ──────────────────────────────────────────────────────────────────────

def main() -> None:
    ap = argparse.ArgumentParser(
        description="Convert any PPTX to ADATA template style (preserves text/images/tables)."
    )
    ap.add_argument("source", type=Path, help="Source PPTX to convert")
    ap.add_argument("output", type=Path, nargs="?",
                    help="Output filename (default: <source>-adata.pptx)")
    args = ap.parse_args()

    if not args.source.exists():
        sys.exit(f"Error: '{args.source}' not found.")
    if not TEMPLATE_PATH.exists():
        sys.exit(f"Error: ADATA template not found at '{TEMPLATE_PATH}'.\n"
                 f"Expected: {TEMPLATE_PATH}")

    out = args.output or args.source.with_stem(args.source.stem + "-adata")
    convert(args.source, out)


if __name__ == "__main__":
    main()
